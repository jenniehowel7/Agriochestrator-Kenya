from datetime import datetime
from pathlib import Path
from typing import TYPE_CHECKING, Any, cast

from pptx import Presentation as create_presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

if TYPE_CHECKING:
    from pptx.presentation import Presentation


TITLE_COLOR = RGBColor(8, 48, 107)
ACCENT_COLOR = RGBColor(0, 123, 85)
BODY_COLOR = RGBColor(33, 37, 41)


def add_title_slide(prs: "Presentation", title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = cast(Any, slide.placeholders[1])
    if title_shape is None:
        return
    title_shape.text = title
    subtitle_shape.text = subtitle
    title_shape.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    subtitle_shape.text_frame.paragraphs[0].font.color.rgb = BODY_COLOR


def add_bullets_slide(prs: "Presentation", title: str, bullets: list[str], footer: str | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    if title_shape is not None:
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR

    body = cast(Any, slide.shapes.placeholders[1]).text_frame
    body.clear()
    for i, text in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = text
        p.level = 0
        p.font.size = Pt(22)
        p.font.color.rgb = BODY_COLOR

    if footer:
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.4))
        tf = tx.text_frame
        tf.text = footer
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.color.rgb = ACCENT_COLOR


def build_deck(output_path: Path) -> None:
    prs = create_presentation()

    add_title_slide(
        prs,
        "AgriOrchestrator Kenya",
        "Autonomous AI copilot for smallholder farm decisions | Hackathon Pitch",
    )

    add_bullets_slide(
        prs,
        "Problem We Are Solving",
        [
            "Kenyan smallholder farmers face weather volatility, pests, and delayed advisory access.",
            "Decisions on irrigation, spraying, and input spend are often made with incomplete data.",
            "Traditional extension support is limited in reach and frequency.",
            "Result: avoidable yield loss, wasted inputs, and unstable household income.",
        ],
        footer="Target users: farmers, extension officers, agri-program admins",
    )

    add_bullets_slide(
        prs,
        "Our Solution",
        [
            "A web platform that combines multi-agent orchestration, predictive risk scoring, and AI advisory.",
            "Farm telemetry + Kenya datasets + model outputs -> clear, actionable daily recommendations.",
            "Role-aware experience for farmer, officer, and admin personas.",
            "Offline-first sync queue keeps operations resilient in low-connectivity contexts.",
        ],
        footer="Core promise: better farm decisions before losses happen",
    )

    add_bullets_slide(
        prs,
        "How It Works (Architecture)",
        [
            "Perception agent: ingests telemetry and context signals.",
            "Forecast agent: estimates near-term risk dynamics.",
            "Reasoning + Action agents: prioritize interventions and generate advisories.",
            "Optimizer agent: improves recommendation strategy over time.",
            "Model Lab: hot-swap trained artifact with visible model card and lineage.",
            "FastAPI mobile backend + Flask frontend + SQLite + M-Pesa simulator.",
        ],
    )

    add_bullets_slide(
        prs,
        "Product Demo Highlights",
        [
            "Premium landing and authentication with role-based navigation.",
            "Dashboard with operational alerts and visual trend summaries.",
            "Profiles management for farmer records and updates.",
            "Payments tab using persistent M-Pesa transaction simulation.",
            "Data Hub and Sync pages for offline queue monitoring and manual flush.",
            "Model Details page exposing version, metrics, SHA256 hash, and training metadata.",
        ],
    )

    add_bullets_slide(
        prs,
        "Data & AI Credibility",
        [
            "Uses Kenya-relevant data sources (NASA weather + World Bank agricultural indicators).",
            "Training notebook supports temporal validation and richer diagnostics.",
            "Model card captures AUC/F1/Brier, features, data sources, and trained years.",
            "Transparent lineage helps trust, auditability, and safer deployment decisions.",
        ],
        footer="This is not a black box demo; it is an inspectable AI system",
    )

    add_bullets_slide(
        prs,
        "Impact, Market, and Execution",
        [
            "Immediate value: lower avoidable losses and better timing of farm interventions.",
            "Adoption path: pilot with county extension teams and farmer groups.",
            "Business path: B2B2C via agri-inputs, NGOs, county programs, and insurers.",
            "Scalability: modular architecture enables additional crops, counties, and channels.",
            "Next 90 days: pilot onboarding, field validation, and model refinement loop.",
        ],
    )

    add_bullets_slide(
        prs,
        "The Ask & Why We Can Win",
        [
            "Ask: support to run pilots, access partners, and accelerate production rollout.",
            "Why now: climate pressure + digital readiness + urgent need for practical farm AI.",
            "Why us: working end-to-end product, real datasets, explainable model layer, and demo-ready UX.",
            "Vision: become the trusted AI decision layer for smallholder agriculture in Kenya.",
        ],
        footer=f"Deck generated on {datetime.now().strftime('%Y-%m-%d')}",
    )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


if __name__ == "__main__":
    output = Path("assets/AgriOrchestrator_Kenya_Pitch_Deck_8_Slides.pptx")
    build_deck(output)
    print(str(output.resolve()))
